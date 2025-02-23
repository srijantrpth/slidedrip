import fitz
from pptx import Presentation
from docx import Document
import os

def extract_text(file_path: str) -> str:
    """Extract Text from PDF, PPTX or DOCX Files """
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        doc = fitz.open(file_path)
        return "\n".join([page.get_text() for page in doc])
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    elif ext==".docx":
        return "\n".join([para.text for para in Document(file_path).paragraphs])
    else:
        return "Unsupported File Format! "